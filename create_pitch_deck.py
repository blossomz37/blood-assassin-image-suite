#!/usr/bin/env python3
"""
Blood Assassin - Story Pitch Deck Generator
Creates a professional PowerPoint presentation for pitching the novel
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from pathlib import Path

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image to slide with error handling"""
    try:
        if os.path.exists(image_path):
            # Check if image needs conversion from WebP
            img = Image.open(image_path)
            if img.format == 'WEBP' or image_path.lower().endswith('.webp'):
                # Convert to PNG
                png_path = image_path.replace('.webp', '.png').replace('.WEBP', '.png')
                if not png_path.endswith('.png'):
                    png_path = image_path + '.png'
                img.save(png_path, 'PNG')
                image_path = png_path
            
            if width and height:
                slide.shapes.add_picture(image_path, left, top, width, height)
            elif width:
                slide.shapes.add_picture(image_path, left, top, width)
            else:
                slide.shapes.add_picture(image_path, left, top)
            return True
    except Exception as e:
        print(f"Could not add image {image_path}: {e}")
    return False

def create_pitch_deck():
    """Create the Blood Assassin pitch deck"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)  # Widescreen
    prs.slide_height = Inches(9)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 25, 35)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "BLOOD ASSASSIN"
    p.font.name = 'Arial Black'
    p.font.size = Pt(72)
    p.font.color.rgb = RGBColor(139, 0, 0)  # Dark red
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "A Dark Fantasy Novel"
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(12), Inches(2))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "When prophecy and blood collide, even monsters must choose their humanity"
    p.font.name = 'Arial Italic'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 170)
    p.alignment = PP_ALIGN.CENTER
    
    # Try to add character portraits if available
    add_image_safe(slide, "generated-images/01_elara_nightshade_portrait.png", 
                   Inches(0.5), Inches(1), height=Inches(7))
    add_image_safe(slide, "generated-images/02_queen_lysandria_portrait.png", 
                   Inches(12), Inches(1), height=Inches(7))
    
    # Slide 2: The World
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "THE WORLD"
    
    # World description
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
    tf = content_box.text_frame
    tf.clear()
    
    # Setting
    p = tf.paragraphs[0]
    p.text = "A REALM DIVIDED"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(139, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "\nThe Crimson Court: Vampire aristocracy ruling from shadow"
    p.font.size = Pt(20)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Order of the Dagger: Elite hunters sworn to protect humanity"
    p.font.size = Pt(20)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Rebellion: Mortals fighting for freedom from supernatural tyranny"
    p.font.size = Pt(20)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nTHE BLOODBOUND PROPHECY"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(139, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "A chosen one with mixed blood will either unite the realms or destroy them all"
    p.font.size = Pt(20)
    p.font.italic = True
    
    # Slide 3: Main Characters
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "CORE CHARACTERS"
    
    # Character grid
    chars = [
        ("ELARA NIGHTSHADE", "The Chosen Hunter", "Torn between humanity and vampiric heritage"),
        ("QUEEN LYSANDRIA", "The Tyrant", "Paranoid ruler obsessed with control"),
        ("SERAPHIEL", "The Fallen Guide", "Celestial being seeking redemption"),
        ("VALERIA DUSKBANE", "The Double Agent", "Spy caught between two worlds")
    ]
    
    y_pos = 2
    for name, role, desc in chars:
        # Name
        name_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(4), Inches(0.5))
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(139, 0, 0)
        
        # Role
        role_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(3), Inches(0.5))
        p = role_box.text_frame.paragraphs[0]
        p.text = role
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(100, 100, 120)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(9), Inches(y_pos), Inches(6), Inches(0.5))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        
        y_pos += 1.5
    
    # Slide 4: The Conflict
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "THE CENTRAL CONFLICT"
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "PERSONAL STAKES"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(139, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "• Elara discovers her vampiric bloodline, shattering her identity"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Must master both sides of her nature to fulfill the prophecy"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Faces betrayal from her mentor who hid the truth"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nGLOBAL STAKES"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(139, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "• The Nightbringer threatens to lock the world in eternal stasis"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Vampire-mortal war escalates toward mutual destruction"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Only Elara can bridge both worlds—or destroy them"
    p.font.size = Pt(18)
    
    # Slide 5: Themes
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "CORE THEMES"
    
    themes = [
        ("Identity & Duality", "The struggle between our nature and our choices"),
        ("Power & Paranoia", "How fear of losing control becomes self-destructive"),
        ("Trust & Betrayal", "Learning vulnerability is strength, not weakness"),
        ("Redemption", "Everyone deserves a chance to change their story"),
        ("Growth vs. Stasis", "The danger of refusing to evolve")
    ]
    
    y_pos = 2
    for theme, desc in themes:
        # Theme name
        theme_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(4), Inches(0.5))
        p = theme_box.text_frame.paragraphs[0]
        p.text = theme
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(139, 0, 0)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(9), Inches(0.5))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.italic = True
        
        y_pos += 1.2
    
    # Slide 6: Target Audience & Comps
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "MARKET POSITIONING"
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "TARGET AUDIENCE"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(139, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "• Adults 25-45 who enjoy dark fantasy and vampire fiction"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Readers of morally complex protagonists and anti-heroes"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Fans of political intrigue and prophecy-driven narratives"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nCOMPARABLE TITLES"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(139, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "• Throne of Glass meets Interview with the Vampire"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• The Witcher's moral complexity with Underworld's aesthetic"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Castlevania's gothic horror meets Game of Thrones' politics"
    p.font.size = Pt(18)
    
    # Slide 7: Series Potential
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "SERIES POTENTIAL"
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "BOOK ONE: Blood Assassin"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(139, 0, 0)
    p = tf.add_paragraph()
    p.text = "Elara discovers her heritage and confronts the Nightbringer"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nBOOK TWO: Crimson Prophecy"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(139, 0, 0)
    p = tf.add_paragraph()
    p.text = "The aftermath of the prophecy creates new alliances and enemies"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nBOOK THREE: Twilight Throne"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(139, 0, 0)
    p = tf.add_paragraph()
    p.text = "Elara must unite both worlds against an ancient threat"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nSPIN-OFF POTENTIAL"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(100, 0, 0)
    p = tf.add_paragraph()
    p.text = "• Seraphiel's celestial past  • Mara's rebellion  • The Order's origins"
    p.font.size = Pt(18)
    
    # Slide 8: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 25, 35)
    
    # Main text
    main_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(5))
    tf = main_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "JOIN THE HUNT"
    p.font.name = 'Arial Black'
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(139, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nA story of monsters, humanity, and the choices that define us"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\n\nBLOOD ASSASSIN"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(180, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Coming Soon"
    p.font.italic = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 150, 170)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "Blood_Assassin_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"✅ Pitch deck created: {output_path}")
    
    # Create a version with placeholder for images if they're not found
    if not os.path.exists("generated-images"):
        print("\n📌 Note: Add character portraits to 'generated-images/' folder for visual enhancement")
    
    return output_path

if __name__ == "__main__":
    # Install required packages if needed
    try:
        from pptx import Presentation
    except ImportError:
        print("Installing python-pptx...")
        import subprocess
        subprocess.run(["pip", "install", "python-pptx", "--break-system-packages"])
        from pptx import Presentation
    
    try:
        from PIL import Image
    except ImportError:
        print("Installing Pillow...")
        import subprocess
        subprocess.run(["pip", "install", "Pillow", "--break-system-packages"])
        from PIL import Image
    
    # Create the pitch deck
    create_pitch_deck()
